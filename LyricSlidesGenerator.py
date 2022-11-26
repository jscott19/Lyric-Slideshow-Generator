# # Install dependencies

# # Install dependencies (python-pptx, lyrics genius) first; run below lines in your terminal (make sure you have PyPI installed so you can use pip)
# pip install python-pptx
# pip install lyricsgenius

# # Define song name and artist
# Change these for each song!

# %%
# Change below lines to whatever the song name and artist are. These should be the only lines you really need to change
SONG_NAME = 'song in my head'
SONG_ARTIST = 'madison cunningham'

# %% [markdown]
# # Setup webscraper to get lyrics

# %%
from lyricsgenius import Genius

token = "MKXDngQJSmRPb_5r87eprRrAdAmfg3RUdduCCPQxwOg35j1idDACH3a7Abtyy-ZQ"
# %env GENIUS_ACCESS_TOKEN=$token
genius = Genius(token)

# %%
artist = genius.search_artist(SONG_ARTIST, max_songs=0)
# print(artist.songs)
song = genius.search_song(SONG_NAME, artist.name)

# %%
song.lyrics

# %% [markdown]
# # Define helpful functions

# %% [markdown]
# ### Parsing Lyrics

# %%
def get_stanza(lyrics):
    # Divide lyrics into lines

    l = len(lyrics)
    # print(l) #debug
    lines = []
    line = ''
    
    i = 0
    in_bracket = in_parens = False
    while i < l:
        char = lyrics[i]
        i += 1
        # if i > l: break
        if char == ']':
            # print('1') #debug
            in_bracket = False
            if char == ']': # stanza complete, next header skipped
                lines.append(line+char)
                line = ''
            continue
        elif char == ')':
            in_parens = False
            continue
        if in_bracket:
            # print('2') #debug
            line += char
            continue
        if in_parens:
            continue # ignore parentheticals
        if char in '[':
            # print('3') #debug
            in_bracket = True
            # line += char
            if line != '': lines.append(line)
            line = char if char == '[' else '' # entering header
        elif char == '(':
            in_parens = True
            continue
        elif char == '\n':
            # print('4') #debug
            if line != '' and line[-1] != char:
                # print('5') #debug
                if line != '': lines.append(line)
                line = ''
        else:
            # print(i) #debug
            line += char
   
    if line != '' and not line.isdigit(): lines.append(line) # append last line


    # Yield stanzas

    isHeader = lambda text: text[0] == '[' and text[-1] == ']'

    j = 1
    stanza = []
    while j < len(lines):
        clean_line = lines[j].strip()  # remove leading or trailing whitespace
        
        # remove end-of-line annotations
        last_word = clean_line.split(' ')[-1]
        if last_word[-1].isdigit() and not last_word.isdigit():
            while last_word[-1].isdigit() and not last_word.isdigit():
                clean_line = clean_line[:-1]
                last_word = clean_line.split(' ')[-1]

        if isHeader(clean_line):
            if stanza != []:
                s = stanza[:]
                stanza = []
                yield s # yield next stanza
        else:
            stanza.append(clean_line)
        j += 1
    yield stanza # yield final stanza

# %% [markdown]
# Test stanza extraction:

# %%
# lyrics = list((song.lyrics.replace('Embed', '')).replace('You might also like', ''))
# stanzas = get_stanza(lyrics)

# # %%
# print(next(stanzas)) #test

# %% [markdown]
# ### Making presentations

# %%
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT

to_rgb = lambda color: RGBColor.from_string(color) if type(color) == str else RGBColor(color[0], color[1], color[2])

def set_bg(slide, color='000000'):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = to_rgb(color)

def add_title_slide(prs, name, artist, color='FFFFFF'):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    set_bg(slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    for box in [title, subtitle]:
        left, top, width, height = Inches(8) - box.width//2, box.top, box.width, box.height
        box.left, box.top, box.width, box.height = left, top, width, height

    title.text = name
    title.text_frame.paragraphs[0].font.color.rgb = to_rgb(color)

    subtitle.text = artist
    

def add_text_slide(prs, text_list, size=40, margin=1, color='FFFFFF'):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide)

    top = width = height = Inches(margin)
    left = Inches(8) - width//2
    top *= 2
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = '\n'.join(text_list)
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for p in tf.paragraphs:
        p.font.color.rgb = to_rgb(color)
        p.font.size = Pt(size)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# %% [markdown]
# # Make the slideshow!

# %%
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(16), Inches(9)

add_title_slide(prs, song.title, artist.name)

# Remove links and annotations from lyrics
lyrics = (song.lyrics.replace('Embed', '')).replace('You might also like', '').split(' ')
for word in lyrics:
    if word[-1].isdigit() and not word.isdigit():
        while word[-1].isdigit() and not word.isdigit():
            word = word[:-1]
lyrics = list(' '.join(lyrics))

for stanza in get_stanza(lyrics): # make lyric slides for each song section (maximum 6 lines per slide)
    num_lines = len(stanza)
    if num_lines < 7:
        add_text_slide(prs, stanza)
    else:
        max_lines_per_slide = 4
        if num_lines % max_lines_per_slide != 0:
            for n in (5,3):
                if num_lines % n == 0:
                    max_lines_per_slide = n
                    break
        i = 0
        while i < num_lines:
            text = stanza[i:i+max_lines_per_slide]
            add_text_slide(prs, text)
            i += max_lines_per_slide

prs.save(song.title+'.pptx')

# %% [markdown]
# Finished! The slideshow should be saved to your local computer as "{Song Title}.pptx" in this folder.


